import os
from pptx import Presentation
from pptx.util import Inches, Pt

def add_title_slide(prs, title, subtitle, team_members):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.title
    subtitle_box = slide.placeholders[1]
    
    title_box.text = title
    
    sub_text = f"{subtitle}\n\nTeam:\n"
    for member in team_members:
        sub_text += f"- {member}\n"
    subtitle_box.text = sub_text

def add_bullet_slide(prs, title, bullets, image_path=None):
    if image_path and os.path.exists(image_path):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = title
        
        tf = body_shape.text_frame
        for i, bullet in enumerate(bullets):
            if i == 0:
                tf.text = bullet
            else:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
                
        # Resize text box to make room for image
        body_shape.width = Inches(4.5)
        # Add image on the right
        slide.shapes.add_picture(image_path, Inches(5.0), Inches(2.0), width=Inches(4.5))
    else:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        shapes = slide.shapes
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = title
        tf = body_shape.text_frame
        for i, bullet in enumerate(bullets):
            if i == 0:
                tf.text = bullet
            else:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0


def create_presentation():
    prs = Presentation()
    
    # Paths
    root_dir = os.path.expanduser('~/sem6/RF/wifi_localization')
    img_dir = os.path.join(root_dir, 'images')
    
    # 1. Title
    add_title_slide(prs, 
        "Real-Time Passive Wi-Fi Zonal Localization",
        "High-Accuracy Spatial Mapping using ESP32 Channel State Information",
        ["<Group Leader Name - Roll No>", "<Partner 2 - Roll No>", "<Partner 3 - Roll No>"])

    # 2. Motivation
    add_bullet_slide(prs, "Motivation & Problem Statement", [
        "Problem: Indoor tracking relies on intrusive cameras or battery-powered wearables.",
        "Privacy: Cameras heavily compromise privacy in vulnerable areas (bedrooms, hospitals).",
        "Solution: Passive Wi-Fi sensing. Humans are mostly water and consistently reflect 2.4GHz RF signals.",
        "We can track human spatial location using ambient Wi-Fi waves already bouncing around the room."
    ])
    
    # 3. What is CSI
    add_bullet_slide(prs, "What is CSI? (Channel State Information)", [
        "Standard Wi-Fi routers use OFDM (Orthogonal Frequency-Division Multiplexing).",
        "The 2.4GHz band is physically split into 64 thinner subcarriers.",
        "CSI gives us the Amplitude and Phase of each individual subcarrier.",
        "This effectively acts as a 64-pixel RF sensor for changes in the environment."
    ])

    # 4. Architecture
    add_bullet_slide(prs, "System Architecture", [
        "Hardware: 1x ESP32 NodeMCU microcontroller.",
        "Capture: Sniffing ambient Wi-Fi packets via custom ESP-IDF firmware.",
        "Processing Unit: Python backend (NumPy, SciPy, Scikit-Learn) attached via Serial over USB.",
        "Output: Zero-latency Live GUI dashboard updating to match human zone localization."
    ])

    # 5. Raw Signal
    add_bullet_slide(prs, "Raw Signal Extraction", [
        "Intercepting string arrays (128 elements) streamed at 115200 baud.",
        "Due to single-antenna clock drift, raw Phase readings are functionally unusable.",
        "We compute Absolute Amplitude via Euclidean distance: Amplitude = sqrt(Real² + Imaginary²).",
        "We discard IEEE 802.11n null subcarriers (27-37), leaving 53 active RF data streams."
    ])

    # 6. DSP Stage 1
    add_bullet_slide(prs, "DSP Pipeline Stage 1: Noise Reduction", [
        "Problem: Environmental interference (Bluetooth, Cosmic rays) causes massive spikes.",
        "Solution - Hampel Filter:",
        "- We apply a rolling median window (size=15).",
        "- Spikes exceeding 3 standard deviations (via Median Absolute Deviation) are replaced.",
        "- This stabilizes the signal for subsequent calculus without dropping frames."
    ])

    # 7. DSP Stage 2 & 3
    add_bullet_slide(prs, "DSP Pipeline Stage 2 & 3: Human Isolation", [
        "Dynamic Background Subtraction: Subtracting a trailing 100-packet moving average zeroes-out the static room geometry.",
        "Butterworth Bandpass Filter: We apply a 3rd-order Butterworth bandpass at [0.1 Hz, 3.0 Hz].",
        "This strictly isolates the Doppler signatures of human breathing (0.1 Hz) and human walking (up to 3.0 Hz)."
    ], os.path.join(img_dir, "dsp_comparison.png"))

    # 8. Data Collection
    add_bullet_slide(prs, "Data Collection Environments", [
        "Environment 1 (Corridor): Zones A & B. Challenging due to perfectly symmetrical multipath geometry.",
        "Environment 2 (Enclosed Room): Empty, Zone A, Zone B, Zone C.",
        "Complex multipath room geometry makes spatial resolution more distinct.",
        "Extracted 1,500 continuous CSI packets (~15 seconds) per zone for model training."
    ], os.path.join(img_dir, "heatmap_room_clean.png"))

    # 9. Feature Engineering
    add_bullet_slide(prs, "Feature Engineering: Curse of Dimensionality", [
        "Standard statistical mapping (Mean, Variance) creates sparse spaces with poor separability.",
        "The Upgrade: We engineered 135 distinct mathematical features focused on Time-Frequency Structure.",
        "Key Feature: Multi-Lag Autocorrelation (Lags 1, 5, 10).",
        "Why: Correlating the window with delayed copies of itself separates erratic noise from steady, rhythmic walking gestures."
    ])

    # 10. ADV Spectral
    add_bullet_slide(prs, "Advanced Spectral & Spatial Features", [
        "Temporal Non-Stationarity: Ratio of variance between halves of a 1-second window detects subjects moving *through* boundaries.",
        "Covariance Eigenvalues: Top-5 Eigenvalues of the subcarrier covariance sparse-matrix.",
        "Why: Eigenvalue clustering maps the complexity of multipath fading. Line-Of-Sight blockages appear distinctly different from flat-fading reflections."
    ], os.path.join(img_dir, "pca_room.png"))

    # 11. ML
    add_bullet_slide(prs, "Machine Learning Selection", [
        "Discarded SVM: Support Vector Machines struggled. They scale poorly in high-dimensional correlation spaces.",
        "Strict StandardScaling ruined our relative subcarrier magnitude physics required by the SVM-RBF kernel.",
        "Chosen Model: Random Forest Ensemble.",
        "Decision trees implicitly feature-select, efficiently ignoring noise and natively managing high-dimensional spatial physics."
    ])

    # 12. Leakage Bug
    add_bullet_slide(prs, "Overcoming The Data Leakage Trap", [
        "The Bug: Initial training reported 96.3% cross-validation accuracy, but totally failed in live tests.",
        "The Cause: A temporal sliding window step of 10 created 90% overlap.",
        "The Result: K-Fold Cross-Validation leaked nearly-identical, correlated frames spanning the training and testing sets.",
        "The model wasn't learning physical zones—it was memorizing random localized noise bumps unique to one specific recording."
    ])

    # 13. Honest Validation
    add_bullet_slide(prs, "Honest Validation & Results", [
        "We corrected the sliding overlap to a strict 50% (step=50) for completely uncorrelated test windows.",
        "Generated 8x training data via Physical Noise Augmentation.",
        "Gaussian noise applied proportionally to per-subcarrier variance mathematically simulated shifting multi-path fading models.",
        "Final Generalized Accuracy: 73.3% across 4 distinct zones. Empty room detection recall exceeded 83%."
    ])

    # 14. Live inference
    add_bullet_slide(prs, "Live Inference Architecture", [
        "Serial I/O bottleneck bypassed using unallocated NumPy 'deque' Ring Buffers for continuous streaming matrix extraction.",
        "UI Stabilization Stage 1: Random Forest must pass a strict >50.0% confidence probability.",
        "UI Stabilization Stage 2: Handled by a 3-vote Exponential Queue.",
        "The 1.0s release timeout prevents the GUI dashboard from flickering wildly, providing a stable user experience."
    ])

    # 15. Conclusion
    add_bullet_slide(prs, "Conclusion & Future Work", [
        "Conclusion: Single-antenna ESP32s can perform high-accuracy spatial mapping and zoning when coupled with heavy DSP and advanced Time-Frequency feature extraction.",
        "Future Integration: Fusing this RF location data with our 'Vinayabrhami' AI OS for autonomous indoor robotics.",
        "Next Steps: Upgrading to multi-antenna MIMO systems.",
        "Enabling Phase-difference processing (CFO/SFO calibration) for direct Angle of Arrival (AoA) tracking protocols."
    ])

    outpath = os.path.join(root_dir, "PRISM_Zonal_Localization.pptx")
    prs.save(outpath)
    print(f"Presentation saved to {outpath}")

if __name__ == '__main__':
    create_presentation()
